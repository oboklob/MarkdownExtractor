from pptx import Presentation
from pptx.util import Pt


def extract_pptx_md(file_path):
    presentation = Presentation(file_path)
    result = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        if run.font.bold:
                            text = f"**{text}**"  # bold text
                        if run.font.italic:
                            text = f"*{text}*"  # italic text
                        if run.font.underline:
                            text = f"_{text}_"  # underline text
                        if run.hyperlink and run.hyperlink.address:
                            text = f"[{text}]({run.hyperlink.address})"  # hyperlink
                        if run.font.size and run.font.size >= Pt(24):
                            text = f"# {text}"  # heading
                        result.append(text)

    return "\n".join(result)
